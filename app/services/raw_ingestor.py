from __future__ import annotations

import hashlib
import sqlite3
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from app.services.ingester import append_log, register_source, stable_source_id_for_rel_path
from app.util.slugs import slugify
from app.util.timestamps import now_iso8601

RAW_SUBDIRS = ("pdfs", "videos", "decks", "audio", "webclips", "spreadsheets", "misc")
TEXT_SUFFIXES = {".txt", ".md", ".markdown", ".csv", ".json", ".yaml", ".yml", ".log"}


@dataclass(slots=True)
class IngestScanSummary:
    scanned_files: int = 0
    processed_files: int = 0
    skipped_files: int = 0
    error_files: int = 0
    rebuilt: bool = False

    def to_dict(self) -> dict[str, Any]:
        return {
            "scanned_files": self.scanned_files,
            "processed_files": self.processed_files,
            "skipped_files": self.skipped_files,
            "error_files": self.error_files,
            "rebuilt": self.rebuilt,
        }


def ingest_db_path(vault_root: Path) -> Path:
    return vault_root / "manifests" / "ingest.sqlite"


def init_ingest_db(vault_root: Path) -> Path:
    db_path = ingest_db_path(vault_root)
    db_path.parent.mkdir(parents=True, exist_ok=True)
    with sqlite3.connect(db_path) as conn:
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS raw_files (
                source_id TEXT PRIMARY KEY,
                rel_path TEXT UNIQUE NOT NULL,
                sha256 TEXT NOT NULL,
                size_bytes INTEGER NOT NULL,
                mtime_ns INTEGER NOT NULL,
                detected_type TEXT NOT NULL,
                status TEXT NOT NULL,
                transcript_rel_path TEXT,
                source_note_rel_path TEXT,
                first_seen TEXT NOT NULL,
                last_seen TEXT NOT NULL,
                last_ingested TEXT,
                last_error TEXT
            )
            """
        )
        conn.execute("CREATE INDEX IF NOT EXISTS idx_raw_files_rel_path ON raw_files(rel_path)")
        conn.commit()
    return db_path


def scan_and_ingest_raw_files(vault_root: Path) -> IngestScanSummary:
    summary = IngestScanSummary()
    db_path = init_ingest_db(vault_root)
    for source_path in iter_raw_files(vault_root):
        summary.scanned_files += 1
        try:
            changed, _ = ingest_single_source(vault_root, source_path, db_path=db_path)
            if changed:
                summary.processed_files += 1
            else:
                summary.skipped_files += 1
        except Exception:  # noqa: BLE001
            summary.error_files += 1
    return summary


def ingest_single_source(vault_root: Path, source_path: Path, *, db_path: Path | None = None) -> tuple[bool, str]:
    """
    Ingest one raw file into vault-local tracking and authoring outputs.
    Returns (changed, source_id) where changed indicates whether processing occurred.
    """
    source_path = source_path.resolve()
    vault_root = vault_root.resolve()
    raw_root = vault_root / "raw"
    if not source_path.exists() or not source_path.is_file():
        raise FileNotFoundError(f"Source file not found: {source_path}")
    if raw_root not in source_path.parents:
        raise ValueError(f"Source file must be inside vault raw/ directory: {source_path}")

    rel = source_path.relative_to(vault_root)
    source_id = stable_source_id_for_rel_path(rel)
    detected_type = detect_type(source_path)
    sha = sha256_file(source_path)
    stat = source_path.stat()
    now = now_iso8601()
    db_path = db_path or init_ingest_db(vault_root)

    with sqlite3.connect(db_path) as conn:
        row = conn.execute(
            "SELECT sha256, size_bytes, mtime_ns FROM raw_files WHERE source_id = ?",
            (source_id,),
        ).fetchone()
        unchanged = bool(row and row[0] == sha and int(row[1]) == stat.st_size and int(row[2]) == stat.st_mtime_ns)
        if unchanged:
            conn.execute(
                """
                UPDATE raw_files
                SET last_seen = ?, status = 'unchanged', last_error = NULL
                WHERE source_id = ?
                """,
                (now, source_id),
            )
            conn.commit()
            return False, source_id

    transcript_rel_path = write_transcript(vault_root, source_id, rel, source_path, detected_type)
    source_note_rel_path = upsert_source_note(vault_root, source_id, rel, detected_type, transcript_rel_path)
    # Keep JSONL manifest as append-only operational audit while SQLite keeps deduped latest state.
    register_source(vault_root, source_path, detected_type=detected_type, source_id=source_id)
    append_log(vault_root, f"Ingested {rel.as_posix()} as {source_id}")

    with sqlite3.connect(db_path) as conn:
        conn.execute(
            """
            INSERT INTO raw_files (
                source_id, rel_path, sha256, size_bytes, mtime_ns, detected_type, status,
                transcript_rel_path, source_note_rel_path, first_seen, last_seen, last_ingested, last_error
            ) VALUES (?, ?, ?, ?, ?, ?, 'ingested', ?, ?, ?, ?, ?, NULL)
            ON CONFLICT(source_id) DO UPDATE SET
                rel_path = excluded.rel_path,
                sha256 = excluded.sha256,
                size_bytes = excluded.size_bytes,
                mtime_ns = excluded.mtime_ns,
                detected_type = excluded.detected_type,
                status = 'ingested',
                transcript_rel_path = excluded.transcript_rel_path,
                source_note_rel_path = excluded.source_note_rel_path,
                last_seen = excluded.last_seen,
                last_ingested = excluded.last_ingested,
                last_error = NULL
            """,
            (
                source_id,
                rel.as_posix(),
                sha,
                stat.st_size,
                stat.st_mtime_ns,
                detected_type,
                transcript_rel_path.as_posix(),
                source_note_rel_path.as_posix(),
                now,
                now,
                now,
            ),
        )
        conn.commit()

    return True, source_id


def iter_raw_files(vault_root: Path) -> list[Path]:
    files: list[Path] = []
    raw_root = vault_root / "raw"
    for sub in RAW_SUBDIRS:
        base = raw_root / sub
        if not base.exists():
            continue
        files.extend([p for p in base.rglob("*") if p.is_file()])
    return sorted(files)


def detect_type(file_path: Path) -> str:
    parent = file_path.parent.name.lower()
    if parent in RAW_SUBDIRS:
        return parent
    suffix = file_path.suffix.lower()
    if suffix in {".pdf"}:
        return "pdfs"
    if suffix in {".ppt", ".pptx"}:
        return "decks"
    if suffix in {".xls", ".xlsx"}:
        return "spreadsheets"
    if suffix in {".mp3", ".wav", ".m4a"}:
        return "audio"
    if suffix in {".mp4", ".mov", ".mkv"}:
        return "videos"
    return "misc"


def sha256_file(file_path: Path) -> str:
    h = hashlib.sha256()
    with file_path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def write_transcript(
    vault_root: Path,
    source_id: str,
    rel_path: Path,
    source_path: Path,
    detected_type: str,
) -> Path:
    text = extract_text(source_path)
    transcript_path = vault_root / "exports" / "transcripts" / f"{source_id}.md"
    transcript_path.parent.mkdir(parents=True, exist_ok=True)
    transcript = (
        f"# Transcript: {source_id}\n\n"
        f"- Source: `{rel_path.as_posix()}`\n"
        f"- Type: `{detected_type}`\n"
        f"- Extracted: `{now_iso8601()}`\n\n"
        "## Content\n\n"
        f"{text}\n"
    )
    transcript_path.write_text(transcript, encoding="utf-8")
    return transcript_path.relative_to(vault_root)


def upsert_source_note(
    vault_root: Path,
    source_id: str,
    rel_path: Path,
    detected_type: str,
    transcript_rel_path: Path,
) -> Path:
    note_slug = slugify(f"{rel_path.stem}-{source_id}") or source_id
    note_rel_path = Path("content/source-notes") / f"{note_slug}.md"
    note_path = vault_root / note_rel_path
    note_path.parent.mkdir(parents=True, exist_ok=True)
    now = now_iso8601()
    frontmatter = (
        "---\n"
        f'id: "{source_id}"\n'
        f'title: "Source Note: {rel_path.name}"\n'
        'type: "source-note"\n'
        'status: "active"\n'
        f'created: "{now}"\n'
        f'updated: "{now}"\n'
        "aliases: []\n"
        "tags: []\n"
        "projects: []\n"
        f'sources: ["{source_id}"]\n'
        "related: []\n"
        "publish: true\n"
        "---\n"
    )
    body = (
        "\n"
        "## Source metadata\n\n"
        f"- Source ID: `{source_id}`\n"
        f"- Raw path: `{rel_path.as_posix()}`\n"
        f"- Detected type: `{detected_type}`\n"
        f"- Transcript: `{transcript_rel_path.as_posix()}`\n\n"
        "## Summary\n\n"
        "_Add a human or LLM-authored summary here._\n\n"
        "## Key points\n\n"
        "- _point 1_\n"
        "- _point 2_\n"
    )
    note_path.write_text(frontmatter + body, encoding="utf-8")
    return note_rel_path


def extract_text(source_path: Path) -> str:
    suffix = source_path.suffix.lower()
    if suffix in TEXT_SUFFIXES:
        return source_path.read_text(encoding="utf-8", errors="replace")
    if suffix == ".pdf":
        return extract_pdf_text(source_path)
    if suffix == ".docx":
        return extract_docx_text(source_path)
    if suffix == ".pptx":
        return extract_pptx_text(source_path)
    return (
        f"_No extractor configured for `{source_path.suffix}` yet._\n\n"
        "The file has been tracked in SQLite and manifests; add extraction support as needed."
    )


def extract_pdf_text(source_path: Path) -> str:
    try:
        from pypdf import PdfReader  # type: ignore
    except Exception:
        return "_PDF extractor unavailable (`pypdf` not installed)._"
    reader = PdfReader(str(source_path))
    chunks: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        chunks.append(f"\n## Page {i}\n\n{text.strip()}\n")
    return "\n".join(chunks).strip() or "_PDF had no extractable text._"


def extract_docx_text(source_path: Path) -> str:
    try:
        import docx  # type: ignore
    except Exception:
        return "_DOCX extractor unavailable (`python-docx` not installed)._"
    doc = docx.Document(str(source_path))
    text = "\n".join(p.text for p in doc.paragraphs if p.text)
    return text.strip() or "_DOCX had no extractable paragraph text._"


def extract_pptx_text(source_path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        return "_PPTX extractor unavailable (`python-pptx` not installed)._"
    prs = Presentation(str(source_path))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                lines.append(text)
        if lines:
            parts.append(f"## Slide {idx}\n\n" + "\n".join(lines))
    return "\n\n".join(parts).strip() or "_PPTX had no extractable text._"


def db_stats(vault_root: Path) -> dict[str, int]:
    db_path = init_ingest_db(vault_root)
    with sqlite3.connect(db_path) as conn:
        total = int(conn.execute("SELECT COUNT(*) FROM raw_files").fetchone()[0])
        ingested = int(conn.execute("SELECT COUNT(*) FROM raw_files WHERE status = 'ingested'").fetchone()[0])
        unchanged = int(conn.execute("SELECT COUNT(*) FROM raw_files WHERE status = 'unchanged'").fetchone()[0])
    return {"total": total, "ingested": ingested, "unchanged": unchanged}
